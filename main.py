import os
import time
from typing import List, Optional
from datetime import datetime, timedelta, timezone  
import re
import tempfile
import asyncio
import json
import logging
import traceback
import uuid

# Web Framework
from fastapi import FastAPI, UploadFile, File, Query, Form, BackgroundTasks, Response, Depends, HTTPException, status, Request
from docx import Document
from pptx import Presentation
from fastapi.responses import JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm
from pydantic import BaseModel
import httpx 
from apscheduler.schedulers.background import BackgroundScheduler

# Security
from passlib.context import CryptContext
from jose import JWTError, jwt
from google.oauth2 import id_token
from google.auth.transport import requests as google_requests
from bakong_khqr import KHQR

# Database
from motor.motor_asyncio import AsyncIOMotorClient
from bson import ObjectId

# Cloud Services
from dotenv import load_dotenv
import cloudinary
import cloudinary.uploader
from google import genai
from google.genai import types

# --- 1. LOGGING & CONFIGURATION ---
load_dotenv()

# Configure Logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[logging.StreamHandler()]
)
logger = logging.getLogger(__name__)

# Validate API Keys
GOOGLE_CLIENT_ID = os.getenv("GOOGLE_CLIENT_ID")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
if not GEMINI_API_KEY:
    logger.error("No GEMINI_API_KEY found. Application may not function correctly.")

client = genai.Client(api_key=GEMINI_API_KEY)

cloudinary.config( 
  cloud_name = os.getenv("CLOUDINARY_CLOUD_NAME"), 
  api_key = os.getenv("CLOUDINARY_API_KEY"), 
  api_secret = os.getenv("CLOUDINARY_SECRET"),
  secure = True
)

# Database Setup
MONGO_URL = os.getenv("MONGO_URL")
mongo_client = AsyncIOMotorClient(MONGO_URL)
db = mongo_client.cv_tracking_db
collection = db.candidates
users_collection = db["users"]
transactions_collection = db["transactions"]
folders_collection = db["folders"]

# Payment Setup
BAKONG_DEV_TOKEN = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJkYXRhIjp7ImlkIjoiZmU2MzNjNDdlOGZlNDQ0YSJ9LCJpYXQiOjE3NjQ5NTIyNDksImV4cCI6MTc3MjcyODI0OX0.tbhgtVlzNrTGhD0mKkN33BgopmENupueM7qa9DsDxOI"
BAKONG_ACCOUNT_ID = "say_vathanak@aclb"
MERCHANT_NAME = "SAY SAKSOVATHANAK"
MERCHANT_CITY = "Phnom Penh"
khqr = KHQR(BAKONG_DEV_TOKEN) 

# Auth Setup
SECRET_KEY = os.getenv("SECRET_KEY", "YOUR_SUPER_SECRET_KEY_CHANGE_THIS_IN_PROD")
ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 60 * 24  # 24 hours
pwd_context = CryptContext(schemes=["bcrypt"], deprecated=["auto"])
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="token")

# --- TELEGRAM CONFIGURATION ---
TELEGRAM_BOT_TOKEN = os.getenv("TELEGRAM_BOT_TOKEN")
TELEGRAM_API_URL = f"https://api.telegram.org/bot{TELEGRAM_BOT_TOKEN}"

if not TELEGRAM_BOT_TOKEN:
    logger.warning("No TELEGRAM_BOT_TOKEN found. Telegram integration disabled.")


# --- 2. BACKGROUND TASKS ---
async def run_async_cleanup():
    logger.info("[Auto-Cleanup] Scanning for files to delete...")
    async for user in users_collection.find({"settings.autoDelete": True}):
        try:
            username = user.get("username")
            settings = user.get("settings", {})
            try:
                days = int(settings.get("retention", "30"))
            except:
                days = 30
                
            cutoff_time = datetime.now() - timedelta(days=days)
            cursor = collection.find({
                "uploaded_by": username,
                "upload_date": {"$lt": cutoff_time.isoformat()},
                "cv_url": {"$ne": None}, 
                "file_status": {"$ne": "Expired"},
                "locked": {"$ne": True}
            })

            async for candidate in cursor:
                url = candidate.get("cv_url")
                if not url or "http" not in url: continue 
                try:
                    # FIX: Re-added [0] to prevent passing a list to Cloudinary
                    public_id = url.split('/')[-1].split('.')[0]
                    logger.info(f" -> Deleting PDF for user {username}: {candidate.get('Name')}")
                    cloudinary.uploader.destroy(public_id)
                    await collection.update_one(
                        {"_id": candidate["_id"]},
                        {"$set": {"cv_url": None, "file_status": "Expired"}}
                    )
                except Exception as e:
                    logger.error(f"Error cleaning {candidate.get('_id')}: {e}")
        except Exception as e:
            logger.error(f"Error processing user {user.get('username')}: {e}")

def extract_text_from_office(file_path: str, ext: str) -> str:
    text_content = []
    try:
        if ext == ".docx":
            doc = Document(file_path)
            for para in doc.paragraphs:
                if para.text.strip(): text_content.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip(): text_content.append(cell.text)
        elif ext == ".pptx":
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip(): text_content.append(shape.text)
    except Exception as e:
        logger.error(f"Text Extraction Error: {e}")
        return ""
    return "\n".join(text_content)

async def process_cv_background(file_content: bytes, filename: str, cv_url: str, candidate_id: str, mime_type: str):
    temp_path = None
    gemini_file = None
    try:
        logger.info(f"[{filename}] 🚀 Started AI processing...")
        suffix = ".bin"
        if "pdf" in mime_type: suffix = ".pdf"
        elif "image" in mime_type: suffix = ".jpg"
        elif "word" in mime_type or filename.endswith(".docx"): suffix = ".docx"
        elif "presentation" in mime_type or filename.endswith(".pptx"): suffix = ".pptx"

        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(file_content)
            temp_path = tmp.name

        prompt = """
        You are an expert HR Data Analyst specializing in the Cambodian labor market. 
        Your task is to extract data from CVs that may have poor spelling, no spacing, or mixed Khmer/English text.

        ### 1. ADDRESS NORMALIZATION & CORRECTION (CRITICAL)
        Cambodia has 25 Provinces, 209 Districts (Khans), and 1,646 Communes (Sangkats). 
        - **Correction:** If the candidate writes "beung salang" or "Beoung Salang", you must correct it to the official standard: "Sangkat Boeung Salang, Khan Tuol Kouk, Phnom Penh".
        - **Structure:** Always return addresses in this exact order: "Sangkat [Name], Khan/Srok [Name], [City/Province]".
        - **No Spacing Fix:** If text is "stengmeanchey", correct it to "Sangkat Steung Meanchey".
        - **Common Mappings:**
            - "toul kork" -> "Khan Tuol Kouk"
            - "chamkarmon" -> "Khan Chamkar Mon"
            - "sen sok" -> "Khan Sen Sok"
            - "beoung keng kang" -> "Sangkat Boeung Keng Kang"
        - **Inference:** If only a Sangkat is mentioned, infer the Khan and Province (e.g., "Phsar Thmey" is in "Khan Daun Penh, Phnom Penh").

        ### 2. DATA EXTRACTION RULES
        - **Name:** Romanize Khmer names (e.g., "សខា" -> "Sokha"). Use Title Case. Fix messy casing (e.g., "sOkHa dARA" -> "Sokha Dara").
        - **Tel:** Standardize to "0xx xxx xxx". Remove "+855". If the user wrote "012-33-44-55", convert to "012 334 455".
        - **School:** Use official names or well-known acronyms (RUPP, ITC, NUM, RULE, SETEC). Correct misspellings (e.g., "RUP" -> "RUPP").
        - **Experience:** Summarize as "Role at Company". If the writing is bad (e.g., "i work at coffee shop before for 2 year"), clean it to "Barista at Coffee Shop".
        - **EducationLevel:** Must be one of: ['High School', 'Associate', 'Bachelor', 'Master', 'PhD', 'Other'].

        ### 3. HANDLING "BAD" WRITING
        If the CV has poor grammar or spelling, use your knowledge of the Cambodian context to "best-guess" the intended meaning so the data matches other high-quality candidates. 

        Return strictly this JSON structure:
        {
            "Name": "Standardized Name",
            "Tel": "0xx xxx xxx",
            "Location": "Sangkat ..., Khan ..., City/Province",
            "School": "Standardized School Name",
            "EducationLevel": "Level",
            "Experience": "Clean Role at Company",
            "Gender": "Male/Female/N/A",
            "BirthDate": "DD-Mon-YYYY",
            "Position": "Target Job Title",
            "Confidence": Integer (0-100)
        }
        """
        is_office_doc = suffix in [".docx", ".pptx"]
        final_prompt = []
        if is_office_doc:
            extracted_text = extract_text_from_office(temp_path, suffix)
            if not extracted_text.strip(): raise Exception("Could not extract text from Office document.")
            final_prompt = [f"{prompt}\n\nHere is the text extracted from the CV:\n{extracted_text}"]
        else:
            gemini_file = client.files.upload(file=temp_path, config=types.UploadFileConfig(mime_type=mime_type))
            while gemini_file.state.name == "PROCESSING":
                await asyncio.sleep(1)
                gemini_file = client.files.get(name=gemini_file.name)
            if gemini_file.state.name == "FAILED": raise Exception("Gemini failed to process the file media.")
            final_prompt = [gemini_file, prompt]

        response = client.models.generate_content(
            model='gemini-2.0-flash', contents=final_prompt,
            config=types.GenerateContentConfig(response_mime_type='application/json')
        )
        try:
            json_text = response.text.replace("```json", "").replace("```", "").strip()
            data = json.loads(json_text)
            # FIX: Added [0] back so lists are extracted properly
            if isinstance(data, list): data = data[0] if len(data) > 0 else {}
        except Exception as parse_error:
            logger.warning(f"[{filename}] JSON Parse Error: {parse_error}")
            data = {"Name": "Parse Error", "Confidence": 0, "Experience": "AI output invalid JSON"}

        update_payload = {
            "Name": data.get("Name", "N/A"), "Tel": data.get("Tel", "N/A"), "Location": data.get("Location", "N/A"),
            "School": data.get("School", "N/A"), "EducationLevel": data.get("EducationLevel", "Other"), 
            "Experience": data.get("Experience", "N/A"), "Gender": data.get("Gender", "N/A"),
            "BirthDate": data.get("BirthDate", "N/A"), "Position": data.get("Position", "N/A"),
            "Confidence": data.get("Confidence", 0), "status": "Ready", "last_modified": datetime.now().isoformat()
        }
        await collection.update_one({"_id": ObjectId(candidate_id)}, {"$set": update_payload})
        logger.info(f"[{filename}] ✅ Success! Confidence: {data.get('Confidence')}")

    except Exception as e:
        logger.error(f"[{filename}] ❌ Background Processing Failed: {str(e)}")
        # FIX: Added traceback logging back
        logger.error(traceback.format_exc())
        user_msg = "Failed to analyze CV."
        if "429" in str(e): user_msg = "System is busy (AI Quota). Retry later."
        elif "json" in str(e).lower(): user_msg = "AI could not read document format."
        await collection.update_one({"_id": ObjectId(candidate_id)}, {"$set": {"status": "Error", "error_msg": user_msg}})
    finally:
        if temp_path and os.path.exists(temp_path):
            try: os.remove(temp_path)
            except: pass
        if gemini_file:
            try: client.files.delete(name=gemini_file.name)
            except: pass

async def process_telegram_cv_pipeline(chat_id: int, user: dict, file_id: str, original_filename: str, mime_type: str):
    """Handles the heavy lifting for Telegram uploads so the webhook can return instantly."""
    try:
        file_content, downloaded_filename = await download_telegram_file(file_id)
        if not file_content:
            await send_telegram_message(chat_id, "❌ Failed to download file from Telegram.")
            return

        filename = original_filename if original_filename != "telegram_upload" else downloaded_filename

        # Deduct Credit
        await users_collection.update_one({"_id": user["_id"]}, {"$inc": {"current_credits": -1, "lifetime_uploads": 1}})
        await transactions_collection.insert_one({
            "user_id": user["_id"], "amount": -1, "type": "SPEND",
            "description": f"Telegram Upload: {filename}", 
            "created_at": datetime.now().isoformat(), "status": "COMPLETED"
        })

        # Upload to Cloudinary
        # FIX: Added [0] so regex works correctly on a string, not a list
        clean_name = re.sub(r'[^a-zA-Z0-9]', '_', filename.split('.')[0]) + "_" + str(uuid.uuid4())[:8]
        upload_result = cloudinary.uploader.upload(file_content, resource_type="auto", public_id=clean_name)
        cv_url = upload_result.get("secure_url")

        # Create DB Placeholder
        placeholder_data = {
            "Name": "Processing...", "Tel": "...", "Location": "...", "School": "...", 
            "Experience": "AI is analyzing...", "Gender": "...", "BirthDate": "...", "Position": "...",
            "file_name": filename, "cv_url": cv_url, "upload_date": datetime.now().isoformat(),
            "locked": False, "status": "Processing", "uploaded_by": user["username"],
            "folder_id": None 
        }
        insert_result = await collection.insert_one(placeholder_data)
        candidate_id = str(insert_result.inserted_id)

        # Send to Gemini processing
        await process_cv_background(file_content, filename, cv_url, candidate_id, mime_type)
        await send_telegram_message(chat_id, f"✅ Success! Document '{filename}' is being parsed. Check your web dashboard shortly.")

    except Exception as e:
        logger.error(f"Telegram Pipeline Error: {e}", exc_info=True)
        await send_telegram_message(chat_id, "❌ An error occurred while processing your file. Credits may need to be refunded.")

# --- 3. APP INITIALIZATION & GLOBAL HANDLERS ---
app = FastAPI()

ADMIN_EMAILS = ["saksovathanaksay@gmail.com"]

@app.exception_handler(HTTPException)
async def http_exception_handler(request: Request, exc: HTTPException):
    return JSONResponse(status_code=exc.status_code, content={"status": "error", "message": exc.detail})

@app.exception_handler(Exception)
async def global_exception_handler(request: Request, exc: Exception):
    error_id = uuid.uuid4()
    logger.error(f"--- Unhandled Error ID: {error_id} | Path: {request.url.path} ---")
    logger.error(f"Details: {str(exc)}")
    # FIX: Added traceback logging back
    logger.error(traceback.format_exc())
    return JSONResponse(status_code=500, content={"status": "error", "message": "An unexpected system error occurred.", "error_id": str(error_id)})

ALLOWED_ORIGINS = ["http://localhost:5173", "https://cvtracker-kh.vercel.app"]
app.add_middleware(
    CORSMiddleware, allow_origins=ALLOWED_ORIGINS, allow_credentials=True,
    allow_methods=["*"], allow_headers=["*"],
)

# --- 4. DATA MODELS ---
class UserCreate(BaseModel):
    username: str
    password: str

class Token(BaseModel):
    access_token: str
    token_type: str

class BulkDeleteRequest(BaseModel):
    candidate_ids: List[str] = []
    
class GoogleAuthRequest(BaseModel):
    token: str
    
class PaymentRequest(BaseModel):
    package_id: str
    email: str
    
class PaymentProof(BaseModel):
    transaction_id: str
    
class UserSettings(BaseModel):
    autoDelete: bool = False
    retention: str = "30"
    exportFields: dict = {}
    autoTags: str = ""
    profile: dict = {}

class FolderCreate(BaseModel): 
    name: str

class MoveRequest(BaseModel):
    candidate_ids: List[str]
    folder_id: Optional[str] = None 

# --- 5. HELPER FUNCTIONS ---
def verify_password(plain_password, hashed_password): return pwd_context.verify(plain_password, hashed_password)
def get_password_hash(password): return pwd_context.hash(password)

def create_access_token(data: dict, expires_delta: Optional[timedelta] = None):
    to_encode = data.copy()
    expire = datetime.now(timezone.utc) + (expires_delta or timedelta(minutes=15))
    to_encode.update({"exp": expire})
    return jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)

async def get_current_user(token: str = Depends(oauth2_scheme)):
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        username: str = payload.get("sub")
        if username is None: raise HTTPException(status_code=401, detail="Invalid token")
    except JWTError: raise HTTPException(status_code=401, detail="Invalid token")
        
    user = await users_collection.find_one({"username": username})
    if user is None: raise HTTPException(status_code=401, detail="User not found")
    return user

async def add_credits(user_id, amount: int, reason: str, ref: str = None):
    await transactions_collection.insert_one({
        "user_id": user_id, "amount": amount, "type": "PURCHASE" if amount > 0 else "SPEND",
        "description": reason, "payment_ref": ref, "created_at": datetime.now().isoformat(), "status": "COMPLETED"
    })
    await users_collection.update_one({"_id": user_id}, {"$inc": {"current_credits": amount}})

async def send_telegram_message(chat_id: int, text: str):
    if not TELEGRAM_BOT_TOKEN: return
    async with httpx.AsyncClient() as client:
        await client.post(f"{TELEGRAM_API_URL}/sendMessage", json={"chat_id": chat_id, "text": text})

async def download_telegram_file(file_id: str):
    async with httpx.AsyncClient() as client:
        file_info_res = await client.get(f"{TELEGRAM_API_URL}/getFile?file_id={file_id}")
        file_info = file_info_res.json()
        if not file_info.get("ok"): 
            logger.error(f"Telegram File Info Error: {file_info}")
            return None, None
        file_path = file_info["result"]["file_path"]
        download_url = f"https://api.telegram.org/file/bot{TELEGRAM_BOT_TOKEN}/{file_path}"
        file_res = await client.get(download_url)
        return file_res.content, file_path.split("/")[-1]

# --- 6. API ENDPOINTS ---

@app.on_event("startup")
async def start_scheduler():
    logger.info("--> System Startup: Checking for expired files...")
    await run_async_cleanup()
    async for user in users_collection.find({"lifetime_uploads": {"$exists": False}}):
        current_count = await collection.count_documents({"uploaded_by": user["username"]})
        await users_collection.update_one({"_id": user["_id"]}, {"$set": {"lifetime_uploads": current_count}})

@app.post("/register", status_code=201)
async def register(user: UserCreate):
    if await users_collection.find_one({"username": user.username}):
        raise HTTPException(status_code=400, detail="Username already registered")
    new_user = {
        "username": user.username, "hashed_password": get_password_hash(user.password),
        "current_credits": 0, "lifetime_spend": 0, "created_at": datetime.now().isoformat(), "provider": "local"
    }
    result = await users_collection.insert_one(new_user)
    await add_credits(result.inserted_id, 10, "Welcome Gift: Free 10 Credits", "SIGNUP_BONUS")
    return {"message": "User created successfully"}

@app.post("/token", response_model=Token)
async def login(form_data: OAuth2PasswordRequestForm = Depends()):
    user = await users_collection.find_one({"username": form_data.username})
    if not user or not verify_password(form_data.password, user["hashed_password"]):
        raise HTTPException(status_code=401, detail="Incorrect username or password")
    token = create_access_token(data={"sub": user["username"]}, expires_delta=timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES))
    return {"access_token": token, "token_type": "bearer"}

@app.post("/auth/google")
async def google_login(request: GoogleAuthRequest):
    try:
        id_info = id_token.verify_oauth2_token(request.token, google_requests.Request(), GOOGLE_CLIENT_ID)
        email = id_info.get("email")
        user = await users_collection.find_one({"username": email})
        if not user:
            result = await users_collection.insert_one({
                "username": email, "hashed_password": "GOOGLE_AUTH_USER", "provider": "google",
                "created_at": datetime.now().isoformat(), "current_credits": 0, "lifetime_spend": 0
            })
            await add_credits(result.inserted_id, 10, "Welcome Gift", "SIGNUP_BONUS")
        token = create_access_token(data={"sub": email}, expires_delta=timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES))
        return {"access_token": token, "token_type": "bearer", "username": email}
    except Exception as e:
        raise HTTPException(status_code=400, detail="Google Authentication failed")

@app.get("/users/me")
async def read_users_me(current_user: dict = Depends(get_current_user)):
    return {"id": str(current_user["_id"]), "username": current_user["username"], "current_credits": current_user.get("current_credits", 0), "settings": current_user.get("settings", {})}

@app.post("/folders")
async def create_folder(folder: FolderCreate, current_user: dict = Depends(get_current_user)):
    new_folder = {"user_id": current_user["_id"], "name": folder.name, "created_at": datetime.now().isoformat()}
    result = await folders_collection.insert_one(new_folder)
    return {"status": "success", "id": str(result.inserted_id), "name": folder.name}

@app.get("/folders")
async def get_folders(current_user: dict = Depends(get_current_user)):
    cursor = folders_collection.find({"user_id": current_user["_id"]}).sort("created_at", -1)
    folders = []
    async for folder in cursor:
        count = await collection.count_documents({"uploaded_by": current_user["username"], "folder_id": str(folder["_id"])})
        folders.append({"id": str(folder["_id"]), "name": folder["name"], "count": count, "created_at": folder["created_at"]})
    return folders

@app.delete("/folders/{folder_id}")
async def delete_folder(folder_id: str, current_user: dict = Depends(get_current_user)):
    folder = await folders_collection.find_one({"_id": ObjectId(folder_id), "user_id": current_user["_id"]})
    if not folder: raise HTTPException(status_code=404, detail="Folder not found")
    await folders_collection.delete_one({"_id": ObjectId(folder_id)})
    await collection.delete_many({"folder_id": folder_id, "uploaded_by": current_user["username"]})
    return {"status": "success", "message": "Folder and contents deleted"}

@app.post("/candidates/move")
async def move_candidates(request: MoveRequest, current_user: dict = Depends(get_current_user)):
    if request.folder_id:
        folder = await folders_collection.find_one({"_id": ObjectId(request.folder_id), "user_id": current_user["_id"]})
        if not folder: raise HTTPException(status_code=404, detail="Target folder not found")
    result = await collection.update_many(
        {"_id": {"$in": [ObjectId(cid) for cid in request.candidate_ids]}, "uploaded_by": current_user["username"]},
        {"$set": {"folder_id": request.folder_id}}
    )
    return {"status": "success", "modified": result.modified_count}

@app.post("/upload-cv")
async def upload_cv(background_tasks: BackgroundTasks, files: List[UploadFile] = File(...), folder_id: str = Form(None), current_user: dict = Depends(get_current_user)):
    results = []
    cost = len(files)
    user_credits = current_user.get("current_credits", 0)
    
    if user_credits < cost:
        raise HTTPException(status_code=402, detail=f"Your credit balance is too low to upload {cost} files. Please top up your account to continue.")

    await users_collection.update_one({"_id": current_user["_id"]}, {"$inc": {"current_credits": -cost, "lifetime_uploads": cost}})
    await transactions_collection.insert_one({
        "user_id": current_user["_id"], "amount": -cost, "type": "SPEND", "description": f"Uploaded {cost} CV(s)", "created_at": datetime.now().isoformat(), "status": "COMPLETED"
    })

    final_folder_id = folder_id if folder_id and folder_id not in ["all", "null"] else None

    for file in files:
        try:
            content = await file.read()
            # FIX: Added [0] to correctly generate the regex string
            clean_name = re.sub(r'[^a-zA-Z0-9]', '_', file.filename.split('.')[0])
            try: upload_result = cloudinary.uploader.upload(content, resource_type="auto", public_id=clean_name)
            except Exception as e: raise Exception("File storage service is unavailable.")
            cv_url = upload_result.get("secure_url")
            
            placeholder_data = {
                "Name": "Processing...", "Tel": "...", "Location": "...", "School": "...", "Experience": "AI is analyzing...", "Gender": "...", "BirthDate": "...", "Position": "...",
                "file_name": file.filename, "cv_url": cv_url, "upload_date": datetime.now().isoformat(), "locked": False, "status": "Processing", "uploaded_by": current_user["username"], "folder_id": final_folder_id 
            }
            insert_result = await collection.insert_one(placeholder_data)
            candidate_id = str(insert_result.inserted_id)
            background_tasks.add_task(process_cv_background, content, file.filename, cv_url, candidate_id, file.content_type)
            placeholder_data["_id"] = candidate_id
            results.append(placeholder_data)
        except Exception as e:
            results.append({"filename": file.filename, "status": "Error", "details": str(e)})
    return {"status": f"Queued {len(results)} files", "details": results, "remaining_credits": user_credits - cost}

@app.post("/webhook/telegram")
async def telegram_webhook(request: Request, background_tasks: BackgroundTasks):
    try:
        payload = await request.json()
        message = payload.get("message")
        if not message: return Response(status_code=200)
            
        chat_id = message["chat"]["id"]
        text = message.get("text", "")
        
        # Link Command
        if text.startswith("/link"):
            parts = text.split(" ")
            if len(parts) == 2:
                username = parts[1].strip()
                user = await users_collection.find_one({"username": username})
                if user:
                    await users_collection.update_one({"_id": user["_id"]}, {"$set": {"telegram_chat_id": str(chat_id)}})
                    await send_telegram_message(chat_id, f"🔗 Success! This chat is now linked to: {username}. You can now forward CVs here.")
                else:
                    await send_telegram_message(chat_id, "❌ Account not found. Check your spelling.")
            else:
                await send_telegram_message(chat_id, "Format: /link your_email@domain.com")
            return Response(status_code=200)

        # Verification
        user = await users_collection.find_one({"telegram_chat_id": str(chat_id)})
        if not user:
            await send_telegram_message(chat_id, "⚠️ Your Telegram is not linked. Type `/link your_email@domain.com` to connect your account.")
            return Response(status_code=200)

        document = message.get("document")
        photo = message.get("photo")
        
        file_id = None
        mime_type = "application/octet-stream"
        original_filename = "telegram_upload"
        
        if document:
            file_id = document["file_id"]
            mime_type = document.get("mime_type", "application/pdf")
            original_filename = document.get("file_name", "document.pdf")
        elif photo:
            photo_obj = photo[-1]
            file_id = photo_obj["file_id"]
            mime_type = "image/jpeg"
            original_filename = f"photo_{file_id}.jpg"
        else:
            await send_telegram_message(chat_id, "👋 Hello! Forward a PDF, Word Doc, or Image here, and I'll send it to your CV Tracker dashboard.")
            return Response(status_code=200)

        if user.get("current_credits", 0) < 1:
            await send_telegram_message(chat_id, "💳 Insufficient credits! Please top up your account on the web platform.")
            return Response(status_code=200)

        await send_telegram_message(chat_id, "⏳ File received! Processing...")
        background_tasks.add_task(process_telegram_cv_pipeline, chat_id, user, file_id, original_filename, mime_type)
        return Response(status_code=200)
        
    except Exception as e:
        logger.error(f"Webhook Error: {e}", exc_info=True)
        return Response(status_code=200)

@app.get("/candidates")
async def get_candidates(page: int = Query(1, ge=1), limit: int = Query(20, le=100), search: str = Query(None), folder_id: str = Query(None), current_user: dict = Depends(get_current_user)):
    query_filter = {"uploaded_by": current_user["username"]}
    if folder_id and folder_id != "all": query_filter["folder_id"] = folder_id
    if search:
        text_regex = {"$regex": search, "$options": "i"}
        clean_digits = re.sub(r'[^0-9]', '', search)
        phone_regex = {"$regex": r"[\s\-\.]*".join(list(clean_digits)), "$options": "i"} if clean_digits else text_regex
        query_filter["$and"] = [{"$or": [{"Name": text_regex}, {"Tel": phone_regex}, {"Tel": text_regex}, {"School": text_regex}, {"Location": text_regex}, {"Position": text_regex}]}]

    total_count = await collection.count_documents(query_filter)
    cursor = collection.find(query_filter).sort("upload_date", -1).skip((page - 1) * limit).limit(limit)
    candidates = [ {**c, "_id": str(c["_id"])} async for c in cursor ]
    return {"data": candidates, "page": page, "limit": limit, "total": total_count}

@app.get("/cv/{candidate_id}")
async def get_candidate_cv(candidate_id: str):
    try:
        candidate = await collection.find_one({"_id": ObjectId(candidate_id)})
        if not candidate: return Response(content="Not Found", status_code=404)
        async with httpx.AsyncClient() as client: response = await client.get(candidate["cv_url"])
        media_type = "image/jpeg" if any(x in candidate["cv_url"] for x in [".jpg", ".png"]) else "application/pdf"
        return Response(content=response.content, media_type=media_type)
    except Exception as e:
        return Response(content="Server Error", status_code=500)

@app.delete("/candidates/{candidate_id}")
async def delete_candidate(candidate_id: str, current_user: dict = Depends(get_current_user)):
    candidate = await collection.find_one({"_id": ObjectId(candidate_id), "uploaded_by": current_user["username"]})
    if not candidate: raise HTTPException(status_code=404, detail="Candidate not found")
    if candidate.get("locked"): raise HTTPException(status_code=403, detail="Candidate is locked")
    await collection.delete_one({"_id": ObjectId(candidate_id)})
    return {"status": "Deleted successfully"}

@app.put("/candidates/{candidate_id}")
async def update_candidate(candidate_id: str, updated_data: dict, current_user: dict = Depends(get_current_user)):
    if "_id" in updated_data: del updated_data["_id"]
    updated_data["last_modified"] = datetime.now().isoformat()
    res = await collection.update_one({"_id": ObjectId(candidate_id), "uploaded_by": current_user["username"]}, {"$set": updated_data})
    if res.matched_count == 0: raise HTTPException(status_code=404, detail="Candidate not found")
    return {"status": "Updated successfully"}

@app.put("/candidates/{candidate_id}/lock")
async def toggle_lock(candidate_id: str, request: dict, current_user: dict = Depends(get_current_user)):
    res = await collection.update_one({"_id": ObjectId(candidate_id), "uploaded_by": current_user["username"]}, {"$set": {"locked": request.get("locked", False)}})
    if res.matched_count == 0: raise HTTPException(status_code=404, detail="Candidate not found")
    return {"status": "success"}

@app.put("/users/settings")
async def update_user_settings(settings: UserSettings, current_user: dict = Depends(get_current_user)):
    settings_dict = settings.dict()
    await users_collection.update_one({"_id": current_user["_id"]}, {"$set": {"settings": settings_dict}})
    return {"status": "success", "settings": settings_dict}

@app.post("/candidates/bulk-delete")
async def bulk_delete_candidates(request: BulkDeleteRequest, current_user: dict = Depends(get_current_user)):
    user_filter = {"uploaded_by": current_user["username"], "locked": {"$ne": True}}
    if request.candidate_ids: user_filter["_id"] = {"$in": [ObjectId(cid) for cid in request.candidate_ids]}
    result = await collection.delete_many(user_filter)
    return {"status": "success", "message": f"Deleted {result.deleted_count} candidates"}
    
@app.post("/candidates/{candidate_id}/retry")
async def retry_parsing(candidate_id: str, background_tasks: BackgroundTasks, current_user: dict = Depends(get_current_user)):
    candidate = await collection.find_one({"_id": ObjectId(candidate_id), "uploaded_by": current_user["username"]})
    if not candidate: raise HTTPException(status_code=404, detail="Candidate not found")
    async with httpx.AsyncClient() as client: response = await client.get(candidate["cv_url"])
    background_tasks.add_task(process_cv_background, response.content, candidate.get("file_name", "retry"), candidate["cv_url"], candidate_id, "application/pdf")
    await collection.update_one({"_id": ObjectId(candidate_id)}, {"$set": {"status": "Processing"}})
    return {"status": "success"}

@app.post("/api/create-payment")
async def create_khqr_payment(request: PaymentRequest):
    packages = {"mini": {"price": 0.25, "credits": 15}, "standard": {"price": 1.50, "credits": 100}, "max": {"price": 5.00, "credits": 400}}
    if request.package_id not in packages: raise HTTPException(status_code=400, detail="Invalid package selected")
    pkg = packages[request.package_id]
    bill_number = str(uuid.uuid4().int)[:10] 
    try:
        qr = khqr.create_qr(bank_account=BAKONG_ACCOUNT_ID, merchant_name=MERCHANT_NAME, merchant_city=MERCHANT_CITY, amount=pkg["price"], currency="USD", phone_number='85592886006', store_label="CV Credits", bill_number=bill_number, terminal_label="POS-01")
    except Exception as e:
        raise HTTPException(status_code=503, detail="Payment service unavailable")
    md5_hash = khqr.generate_md5(qr)
    user = await users_collection.find_one({"username": request.email})
    if user:
        await transactions_collection.insert_one({"user_id": user["_id"], "amount": pkg["credits"], "price": pkg["price"], "type": "PURCHASE_INTENT", "status": "PENDING", "payment_ref": bill_number, "md5_hash": md5_hash, "created_at": datetime.now().isoformat()})
    return {"qr_code": qr, "md5": md5_hash, "amount": pkg["price"]}

@app.post("/api/check-payment-status")
async def check_payment_status(md5_hash: str, force: bool = Query(False)):
    trx = await transactions_collection.find_one({"md5_hash": md5_hash})
    if not trx: raise HTTPException(status_code=404, detail="Transaction not found")
    if trx.get("status") == "COMPLETED": 
        user = await users_collection.find_one({"_id": trx["user_id"]})
        return {"status": "PAID", "message": "Already processed", "new_credits": trx["amount"], "total_credits": user.get("current_credits", 0)}

    payment_status = "UNPAID"
    try:
        if khqr.check_payment(md5_hash) == "PAID": payment_status = "PAID"
    except Exception as e: pass
    if force: payment_status = "PAID"

    if payment_status == "PAID": 
        await transactions_collection.update_one({"_id": trx["_id"]}, {"$set": {"status": "COMPLETED", "paid_at": datetime.now().isoformat()}})
        await add_credits(trx["user_id"], trx["amount"], f"Purchased Credits (Ref: {trx.get('payment_ref', 'N/A')})", md5_hash)
        updated_user = await users_collection.find_one({"_id": trx["user_id"]})
        return {"status": "PAID", "new_credits": trx["amount"], "total_credits": updated_user.get("current_credits", 0)}
    return {"status": "UNPAID", "detail": "Payment not received yet"}

@app.get("/admin/transactions")
async def get_all_transactions(current_user: dict = Depends(get_current_user)):
    filter_query = {"status": {"$in": ["VERIFYING", "COMPLETED", "REJECTED"]}}
    cursor = transactions_collection.find(filter_query).sort("created_at", -1).limit(50)
    transactions = []
    async for trx in cursor:
        user = await users_collection.find_one({"_id": trx["user_id"]})
        username = user["username"] if user else "Unknown"
        transactions.append({"id": str(trx["_id"]), "username": username, "amount": trx["amount"], "price": trx.get("price", 0), "type": trx["type"], "status": trx["status"], "md5_hash": trx.get("md5_hash"), "payment_ref": trx.get("payment_ref"), "proof_url": trx.get("proof_url"), "created_at": trx["created_at"]})
    return transactions

@app.delete("/admin/transactions")
async def clear_all_transactions(current_user: dict = Depends(get_current_user)):
    await transactions_collection.delete_many({})
    return {"status": "success", "message": "Transaction history wiped."}

@app.post("/admin/add-credits")
async def admin_add_credits(username: str, amount: int):
    user = await users_collection.find_one({"username": username})
    if not user: raise HTTPException(status_code=404, detail="User not found")
    await add_credits(user_id=user["_id"], amount=amount, reason="Admin Manual Top-up", ref="DEV_CHEAT")
    return {"status": "success", "message": f"Added {amount} credits to {username}"}

@app.post("/api/submit-payment-proof")
async def submit_payment_proof(background_tasks: BackgroundTasks, md5_hash: str = Query(...), file: UploadFile = File(...), current_user: dict = Depends(get_current_user)):
    trx = await transactions_collection.find_one({"md5_hash": md5_hash, "user_id": current_user["_id"]})
    if not trx: raise HTTPException(status_code=404, detail="Transaction not found")
    try:
        content = await file.read()
        upload_result = cloudinary.uploader.upload(content, folder="payment_proofs", resource_type="image")
        proof_url = upload_result.get("secure_url")
    except Exception as e: raise HTTPException(status_code=500, detail="Failed to upload proof image")

    await transactions_collection.update_one({"_id": trx["_id"]}, {"$set": {"status": "VERIFYING", "proof_url": proof_url, "submitted_at": datetime.now().isoformat()}})
    return {"status": "success", "message": "Proof submitted for review"}

@app.get("/admin/pending-transactions")
async def get_pending_transactions(current_user: dict = Depends(get_current_user)):
    cursor = transactions_collection.find({"status": "VERIFYING"}).sort("submitted_at", -1)
    results = []
    async for trx in cursor:
        user = await users_collection.find_one({"_id": trx["user_id"]})
        results.append({"id": str(trx["_id"]), "username": user.get("username", "Unknown"), "amount": trx["amount"], "payment_ref": trx.get("payment_ref"), "proof_url": trx.get("proof_url"), "submitted_at": trx.get("submitted_at"), "md5_hash": trx.get("md5_hash")})
    return results

@app.post("/admin/process-transaction/{transaction_id}")
async def process_transaction(transaction_id: str, action: str = Query(..., pattern="^(APPROVE|REJECT)$"), current_user: dict = Depends(get_current_user)):
    trx = await transactions_collection.find_one({"_id": ObjectId(transaction_id)})
    if not trx: raise HTTPException(status_code=404, detail="Transaction not found")
    if trx["status"] == "COMPLETED": return {"status": "error", "message": "Already completed"}

    if action == "APPROVE":
        await add_credits(trx["user_id"], trx["amount"], f"Top-up Approved (Ref: {trx.get('payment_ref')})", trx.get("md5_hash"))
        money_amount = trx.get("price", 0) 
        if money_amount == 0 and trx["amount"] == 20: money_amount = 1.00 
        if money_amount == 0 and trx["amount"] == 150: money_amount = 5.00
        await users_collection.update_one({"_id": trx["user_id"]}, {"$inc": {"lifetime_spend": money_amount}})
        await transactions_collection.update_one({"_id": trx["_id"]}, {"$set": {"status": "COMPLETED", "reviewed_by": current_user["username"], "reviewed_at": datetime.now().isoformat()}})
        return {"status": "success", "message": "Payment Approved & Credits Added"}
    else: 
        await transactions_collection.update_one({"_id": trx["_id"]}, {"$set": {"status": "REJECTED", "reviewed_by": current_user["username"], "reviewed_at": datetime.now().isoformat()}})
        return {"status": "success", "message": "Payment Rejected"}
    
def verify_admin_access(current_user: dict = Depends(get_current_user)):
    user_email = current_user.get("username") 
    if user_email not in ADMIN_EMAILS: raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail="You do not have permission to access the Admin Panel.")
    return current_user

@app.get("/admin/dashboard-stats")
async def get_admin_stats(admin_user: dict = Depends(verify_admin_access)):
    return {"total_candidates": 20, "unique_roles": 4, "message": f"Hello Admin {admin_user.get('username')}, here is your secret data."}

@app.get("/admin/users")
async def get_all_users(current_user: dict = Depends(verify_admin_access)):
    cursor = users_collection.find()
    users = []
    async for user in cursor:
        upload_count = await collection.count_documents({"uploaded_by": user["username"]})
        users.append({"id": str(user["_id"]), "username": user["username"], "credits": user.get("current_credits", 0), "joined_at": user.get("created_at"), "uploads": upload_count, "is_active": user.get("is_active", True)})
    return users

@app.put("/admin/users/{user_id}/toggle-status")
async def toggle_user_status(user_id: str, current_user: dict = Depends(verify_admin_access)):
    user = await users_collection.find_one({"_id": ObjectId(user_id)})
    if not user: raise HTTPException(404, "User not found")
    new_status = not user.get("is_active", True)
    await users_collection.update_one({"_id": ObjectId(user_id)}, {"$set": {"is_active": new_status}})
    return {"status": "success", "is_active": new_status}

@app.delete("/admin/users/{user_id}")
async def delete_user(user_id: str, current_user: dict = Depends(verify_admin_access)):
    user = await users_collection.find_one({"_id": ObjectId(user_id)})
    if not user: raise HTTPException(404, "User not found")
    if user["username"] == current_user["username"]: raise HTTPException(400, "You cannot delete your own admin account.")
    await users_collection.delete_one({"_id": ObjectId(user_id)})
    return {"status": "success", "message": f"User {user['username']} deleted permanently."}

@app.get("/admin/analytics")
async def get_analytics(current_user: dict = Depends(verify_admin_access)):
    pipeline_revenue = [{"$group": {"_id": None, "total_revenue": {"$sum": "$lifetime_spend"}}}]
    revenue_res = await users_collection.aggregate(pipeline_revenue).to_list(length=1)
    total_revenue = revenue_res["total_revenue"] if revenue_res else 0

    pipeline_uploads = [{"$group": {"_id": None, "total_uploads": {"$sum": "$lifetime_uploads"}}}]
    uploads_res = await users_collection.aggregate(pipeline_uploads).to_list(length=1)
    total_files = uploads_res["total_uploads"] if uploads_res else 0
    pending_count = await transactions_collection.count_documents({"status": "VERIFYING"})

    return {"revenue": total_revenue, "pending_reviews": pending_count, "total_files_parsed": total_files}

@app.get("/admin/failed-parsings")
async def get_failed_parsings(current_user: dict = Depends(verify_admin_access)):
    cursor = collection.find({"status": "Error"}).sort("upload_date", -1).limit(50)
    failed_files = []
    async for doc in cursor: failed_files.append({"id": str(doc["_id"]), "file_name": doc.get("file_name"), "uploaded_by": doc.get("uploaded_by"), "error_msg": doc.get("error_msg"), "upload_date": doc.get("upload_date"), "cv_url": doc.get("cv_url")})
    return failed_files